#!/usr/bin/env python
"""Generate the Hierarchical Latent Reasoning Architecture proposal deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------- palette (midnight / teal "latent reasoning" theme) ----------
BG        = "0C1424"   # midnight navy (dominant)
BG2       = "0A101D"   # deeper panel
CARD      = "16233D"   # surface card
CARD2     = "1B2C4A"   # lighter card
LINE      = "26385C"   # hairline / divider
TEAL      = "2DD4BF"   # primary accent
TEALDIM   = "0F766E"
BLUE      = "5B9DF9"   # secondary accent
AMBER     = "F5A524"   # warning accent
CORAL     = "F26D6D"   # danger
GREEN     = "34D399"
TEXT      = "EAF0FA"   # near-white body
MUTE      = "93A7C6"   # muted label
FAINT     = "63769A"   # very muted

HEAD = "Cambria"
BODY = "Calibri"

def C(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

# ---------------- primitives ----------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(bg)
    return s

def _set_font(run, size, color, bold=False, italic=False, font=BODY, spacing=None):
    run.font.size = Pt(size)
    run.font.color.rgb = C(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(spacing*100)))

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space_after=None, line=None):
    """runs: list of paragraphs; each paragraph = list of (txt, opts) run-tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None: p.space_after = Pt(space_after)
        if line is not None: p.line_spacing = line
        for (txt, opts) in para:
            r = p.add_run(); r.text = txt
            _set_font(r, **opts)
    return tb

def rect(s, x, y, w, h, fill=None, lc=None, lw=1.0, radius=0.08,
         shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if lc is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(lc); sp.line.width = Pt(lw)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if not shadow:
        sp.shadow.inherit = False
    else:
        sp.shadow.inherit = False
        _soft_shadow(sp)
    return sp

def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    outer = el.makeelement(qn('a:outerShdw'),
        {'blurRad':'90000','dist':'38100','dir':'5400000','rotWithShape':'0'})
    clr = outer.makeelement(qn('a:srgbClr'), {'val':'000000'})
    alpha = clr.makeelement(qn('a:alpha'), {'val':'42000'})
    clr.append(alpha); outer.append(clr); el.append(outer); spPr.append(el)

def line(s, x1, y1, x2, y2, color=LINE, w=1.0, dash=None):
    ln = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = C(color); ln.line.width = Pt(w)
    if dash:
        d = ln.line._get_or_add_ln()
        pd = d.makeelement(qn('a:prstDash'), {'val': dash})
        d.append(pd)
    ln.shadow.inherit = False
    return ln

def shape_text(sp, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               ml=0.1, mr=0.1, mt=0.05, mb=0.05):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, opts) in para:
            r = p.add_run(); r.text = txt
            _set_font(r, **opts)

def dot(s, x, y, d, color):
    return rect(s, x, y, d, d, fill=color, radius=None, shape=MSO_SHAPE.OVAL)

def arrow(s, x1, y1, x2, y2, color=TEAL, w=1.6, dashed=False, head=0.15):
    """Line from (x1,y1)->(x2,y2) with a triangle head at the end.
    Handles right / left / down / up (the only cases the diagrams use)."""
    line(s, x1, y1, x2, y2, color=color, w=w, dash=('dash' if dashed else None))
    dx, dy = x2 - x1, y2 - y1
    if abs(dx) >= abs(dy):
        rot = 90 if dx >= 0 else 270          # right / left
    else:
        rot = 180 if dy >= 0 else 0           # down / up
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x2-head/2), Inches(y2-head/2), Inches(head), Inches(head))
    tri.rotation = rot
    tri.fill.solid(); tri.fill.fore_color.rgb = C(color)
    tri.line.fill.background(); tri.shadow.inherit = False
    return tri

def eyebrow(s, label, x=0.7, y=0.52, color=TEAL):
    dot(s, x, y+0.02, 0.16, color)
    text(s, x+0.28, y-0.09, 8, 0.4,
         [[(label.upper(), dict(size=13.5, color=color, bold=True, font=BODY, spacing=2.6))]])

def title(s, t, x=0.7, y=0.82, w=12, size=33, color=TEXT):
    text(s, x, y, w, 0.9, [[(t, dict(size=size, color=color, bold=True, font=HEAD))]])

TOTAL = 11
def pagenum(s, n, label):
    text(s, 0.7, 7.02, 6, 0.3,
         [[(label, dict(size=9.5, color=FAINT, font=BODY, spacing=1.5))]])
    text(s, 11.4, 7.02, 1.23, 0.3,
         [[("%02d / %02d" % (n, TOTAL), dict(size=9.5, color=FAINT, font=BODY, spacing=1.5))]],
         align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = slide(BG2)
# faint concentric "thought" rings motif, top-right
for i, (d, a) in enumerate([(6.4,"14000"),(5.0,"20000"),(3.6,"30000"),(2.2,"48000")]):
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(SW-3.0-d/2), Inches(-1.1-d/2+2.0), Inches(d), Inches(d))
    ov.fill.background()
    ov.line.color.rgb = C(TEAL); ov.line.width = Pt(1.1)
    ln = ov.line._get_or_add_ln()
    a_el = ln.makeelement(qn('a:solidFill'), {})
    c_el = a_el.makeelement(qn('a:srgbClr'), {'val': TEAL})
    al = c_el.makeelement(qn('a:alpha'), {'val': a}); c_el.append(al); a_el.append(c_el)
    # replace default solidFill
    for existing in ln.findall(qn('a:solidFill')): ln.remove(existing)
    ln.insert(0, a_el)
    ov.shadow.inherit = False
dot(s, SW-3.0-0.11, 2.0-0.11, 0.22, TEAL)

eyebrow(s, "Project Proposal · 2026", x=0.9, y=1.55, color=TEAL)
text(s, 0.85, 2.05, 11.4, 2.4, [
    [("Hierarchical Latent", dict(size=54, color=TEXT, bold=True, font=HEAD))],
    [("Reasoning Architecture", dict(size=54, color=TEAL, bold=True, font=HEAD))],
], line=1.02)
text(s, 0.9, 4.35, 10.6, 1.0, [
    [("A model that ", dict(size=18, color=MUTE, font=BODY)),
     ("thinks in latent thoughts", dict(size=18, color=TEXT, bold=True, font=BODY)),
     (" — chunk-level vectors deliberated by a", dict(size=18, color=MUTE, font=BODY))],
    [("recurrent loop, then decoded to language by a separate module.",
      dict(size=18, color=MUTE, font=BODY))],
], line=1.28)
# source-paper chips
chips = ["JEPA-Reasoner", "HRM-Text", "Thought Gestalt", "Parcae"]
cx = 0.9
for ch in chips:
    w = 0.34 + len(ch)*0.105
    rect(s, cx, 5.55, w, 0.44, fill=CARD, lc=LINE, lw=1, radius=0.5)
    text(s, cx, 5.55, w, 0.44, [[(ch, dict(size=12.5, color=TEXT, font=BODY))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + 0.22
line(s, 0.9, 6.55, 12.43, 6.55, color=LINE, w=1)
text(s, 0.9, 6.72, 11.5, 0.4,
     [[("Composing four architectures whose problems compose — into one that thinks before it speaks.",
        dict(size=12.5, color=FAINT, italic=True, font=BODY))]])

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
s = slide(BG)
eyebrow(s, "The Problem", color=CORAL)
title(s, "Today's LLMs reason in the same tokens they speak")
text(s, 0.7, 1.62, 12, 0.6, [[
    ("One representation is forced to do two incompatible jobs — and the coupling has three costs.",
     dict(size=15.5, color=MUTE, font=BODY))]])

cards = [
    (CORAL, "Compounding fragility",
     "Autoregression conditions on its own sampled output. One wrong token pollutes the context and corrupts every step of reasoning that follows."),
    (AMBER, "Reasoning tangled with expression",
     "The act of thinking and the act of wording it share the same weights and states — the model cannot deliberate without committing to surface tokens."),
    (BLUE, "No self / input boundary",
     "Nothing structurally separates “the user asserted X” from “I concluded X.” They collapse into one latent — the exact substrate sycophancy exploits."),
]
cw, gap = 3.86, 0.35
cx = 0.7; cy = 2.5; ch = 3.35
for col, head, body in cards:
    rect(s, cx, cy, cw, ch, fill=CARD, lc=LINE, lw=1, radius=0.06, shadow=True)
    rect(s, cx+0.34, cy+0.36, 0.5, 0.5, fill=BG2, lc=col, lw=1.5, radius=0.5)
    dot(s, cx+0.5, cy+0.52, 0.18, col)
    text(s, cx+0.34, cy+1.06, cw-0.68, 0.8,
         [[(head, dict(size=18, color=TEXT, bold=True, font=HEAD))]], line=1.0)
    text(s, cx+0.34, cy+1.95, cw-0.68, 1.3,
         [[(body, dict(size=13, color=MUTE, font=BODY))]], line=1.24)
    cx += cw + gap

text(s, 0.7, 6.15, 12, 0.6, [[
    ("The goal: ", dict(size=14.5, color=TEXT, bold=True, font=BODY)),
    ("decouple reasoning from expression, make deliberation a real compute dial, and give the model a place to hold a belief that isn't just the user's.",
     dict(size=14.5, color=MUTE, font=BODY))]], line=1.2)
pagenum(s, 2, "The Problem")

# ============================================================
# SLIDE 3 — LITERATURE REVIEW
# ============================================================
s = slide(BG)
eyebrow(s, "Literature Review", color=BLUE)
title(s, "Four papers, four different fixes — none combines them")

items = [
    (TEAL, "JEPA-Reasoner",
     "Reasons in latent space; a separate Talker head reads latents out to tokens.",
     "Ablation: the Talker is a pure readout — useless without good latents."),
    (BLUE, "HRM-Text",
     "Dual-timescale recurrence: a fast L-module refines locally, a slow H-module holds context.",
     "Looping beats width, FLOP-for-FLOP — depth as a compute source."),
    (GREEN, "Thought Gestalt",
     "Generates one sentence at a time, cross-attending to a memory of prior sentence vectors.",
     "Gradient flows back through memory — fixes reversal-curse gaps."),
    (AMBER, "Parcae",
     "Stabilizes looped transformers by constraining the spectral norm of the loop's injection.",
     "Test-time looping scales as a predictable, saturating exponential."),
]
gx, gy = 0.7, 1.72
cw, chh, gpx, gpy = 5.86, 1.86, 0.35, 0.28
for i, (col, name, what, result) in enumerate(items):
    x = gx + (i % 2) * (cw + gpx)
    y = gy + (i // 2) * (chh + gpy)
    rect(s, x, y, cw, chh, fill=CARD, lc=LINE, lw=1, radius=0.06, shadow=True)
    rect(s, x+0.3, y+0.32, 0.14, 1.2, fill=col, lc=None, radius=0.5)  # small pill marker (not edge stripe)
    text(s, x+0.62, y+0.28, cw-0.9, 0.4,
         [[(name, dict(size=17, color=col, bold=True, font=HEAD))]])
    text(s, x+0.62, y+0.72, cw-0.9, 0.7,
         [[(what, dict(size=12.5, color=TEXT, font=BODY))]], line=1.16)
    text(s, x+0.62, y+1.42, cw-0.9, 0.4,
         [[("→  ", dict(size=12, color=col, bold=True, font=BODY)),
           (result, dict(size=11.5, color=MUTE, italic=True, font=BODY))]], line=1.1)

# the gap band
gy2 = gy + 2*chh + gpy + 0.16
rect(s, 0.7, gy2, 12.08, 1.02, fill=CARD2, lc=TEAL, lw=1.2, radius=0.06)
dot(s, 1.05, gy2+0.44, 0.2, TEAL)
text(s, 1.42, gy2+0.16, 11.1, 0.8, [[
    ("The gap.  ", dict(size=14.5, color=TEAL, bold=True, font=BODY)),
    ("Each paper fixes one failure in isolation. No prior work composes them — and naively combining latent SSL, un-detached memory, and a looped recurrence creates ",
     dict(size=13.5, color=TEXT, font=BODY)),
    ("new interaction failures", dict(size=13.5, color=TEXT, bold=True, font=BODY)),
    (" that none of the source papers ever had to face.", dict(size=13.5, color=TEXT, font=BODY))]],
    line=1.2, anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 3, "Literature Review")

# ============================================================
# SLIDE 4 — PROPOSED SOLUTION: THE FULL PIPELINE
# ============================================================
s = slide(BG)
eyebrow(s, "Proposed Solution", color=TEAL)
title(s, "The full pipeline: two lanes, one thought loop")
text(s, 0.7, 1.52, 12.0, 0.4, [[
    ("Input enters ", dict(size=13.5, color=MUTE, font=BODY)),
    ("read-only", dict(size=13.5, color=BLUE, bold=True, font=BODY)),
    (" — only the thought loop writes self-state. Each finished thought is remembered, then spoken.",
     dict(size=13.5, color=MUTE, font=BODY))]], line=1.15)

# ---------- INPUT LANE (left column) ----------
ix, iy, iw, ih = 0.7, 2.12, 3.05, 4.24
rect(s, ix, iy, iw, ih, fill=CARD, lc=BLUE, lw=1.3, radius=0.05, shadow=True)
inx, inw = ix+0.25, iw-0.5
text(s, inx, iy+0.2, inw, 0.3, [[("INPUT LANE", dict(size=13, color=BLUE, bold=True, font=BODY, spacing=1.8))]])
text(s, inx, iy+0.5, inw, 0.25, [[("read-only · never writes self-state", dict(size=9.5, color=MUTE, font=BODY))]])
# two source chips
for k,(a,b) in enumerate([("Raw tokens","current turn"),("Aged gestalt summaries","prior turns")]):
    cy = iy+0.92 + k*0.72
    rect(s, inx, cy, inw, 0.6, fill=BG2, lc=LINE, lw=1, radius=0.12)
    text(s, inx+0.16, cy+0.09, inw-0.3, 0.25, [[(a, dict(size=11.5, color=TEXT, bold=True, font=BODY))]])
    text(s, inx+0.16, cy+0.33, inw-0.3, 0.22, [[(b, dict(size=9.5, color=MUTE, font=BODY))]])
arrow(s, ix+iw/2, iy+2.4, ix+iw/2, iy+2.66, color=BLUE, w=1.4)
# input-lane encoder
rect(s, inx, iy+2.72, inw, 0.62, fill=CARD2, lc=BLUE, lw=1.2, radius=0.12)
text(s, inx, iy+2.83, inw, 0.25, [[("Input-lane encoder", dict(size=11.5, color=BLUE, bold=True, font=BODY))]], align=PP_ALIGN.CENTER)
text(s, inx, iy+3.06, inw, 0.22, [[("bidirectional", dict(size=9.5, color=MUTE, font=BODY))]], align=PP_ALIGN.CENTER)
text(s, inx, iy+3.5, inw, 0.7, [
    [("Role tags  ", dict(size=10, color=MUTE, font=BODY)),
     ("USER · SELF · SYSTEM", dict(size=10, color=TEXT, bold=True, font=BODY))],
    [("The loop only cross-attends to it — the substrate for a self / input boundary.",
      dict(size=9.5, color=FAINT, italic=True, font=BODY))],
], line=1.16, space_after=3)

# ---------- GENERATION SPINE (right) ----------
sx, sw = 5.55, 6.95
scx = sx + sw/2
# 1. chunk encoder
y1 = 2.12
rect(s, sx, y1, sw, 0.6, fill=CARD, lc=TEAL, lw=1.3, radius=0.07, shadow=True)
text(s, sx, y1, sw, 0.6, [[
    ("CHUNK ENCODER", dict(size=12, color=TEAL, bold=True, font=BODY, spacing=1.4)),
    ("     chunk cₜ  →  192-d thought latent", dict(size=10.5, color=MUTE, font=BODY))]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, scx, y1+0.6, scx, y1+0.88, color=TEAL)
# 2. inner HRM loop
y2 = 3.02; h2 = 1.2
rect(s, sx, y2, sw, h2, fill=CARD, lc=TEAL, lw=1.7, radius=0.07, shadow=True)
text(s, sx, y2+0.13, sw, 0.28, [[("INNER HRM LOOP  ·  SELF LANE", dict(size=12, color=TEAL, bold=True, font=BODY, spacing=1.2))]], align=PP_ALIGN.CENTER)
lxx = scx-1.06; mmy = y2+0.5
for k in range(3):
    rect(s, lxx, mmy, 0.42, 0.38, fill=BG2, lc=LINE, lw=1, radius=0.16)
    text(s, lxx, mmy, 0.42, 0.38, [[("L", dict(size=12, color=BLUE, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lxx += 0.5
rect(s, lxx+0.12, mmy, 0.5, 0.38, fill=TEALDIM, lc=TEAL, lw=1, radius=0.16)
text(s, lxx+0.12, mmy, 0.5, 0.38, [[("H", dict(size=12, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, sx, y2+0.95, sw, 0.24, [[("fast L ×3  →  slow H ×1   ·   Parcae-stable · adaptive depth · writes self-state",
    dict(size=10, color=MUTE, font=BODY))]], align=PP_ALIGN.CENTER)
arrow(s, scx, y2+h2, scx, y2+h2+0.28, color=TEAL)
text(s, scx+0.15, y2+h2+0.01, 2.2, 0.26, [[("write thought zₜ", dict(size=9.5, color=TEAL, italic=True, font=BODY))]])
# 3. gestalt memory
y3 = 4.5
rect(s, sx, y3, sw, 0.74, fill=CARD2, lc=BLUE, lw=1.3, radius=0.08, shadow=True)
text(s, sx+0.3, y3+0.14, sw-2.4, 0.28, [[("GESTALT MEMORY", dict(size=12, color=BLUE, bold=True, font=BODY, spacing=1.4))]])
text(s, sx+0.3, y3+0.42, sw-2.4, 0.24, [[("FIFO · role-tagged · ungated gradient · next thought recalls it at O(1)", dict(size=10, color=MUTE, font=BODY))]])
zx = sx+sw-2.05
for k in range(3):
    rect(s, zx, y3+0.22, 0.5, 0.34, fill=BG2, lc=BLUE, lw=1, radius=0.16)
    text(s, zx, y3+0.22, 0.5, 0.34, [[("z%d"%(k+1), dict(size=10, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    zx += 0.56
text(s, zx-0.02, y3+0.26, 0.4, 0.3, [[("…", dict(size=13, color=MUTE, font=BODY))]])
arrow(s, scx, y3+0.74, scx, y3+1.02, color=GREEN)
# 4. talker
y4 = 5.54
rect(s, sx, y4, sw, 0.72, fill=CARD, lc=GREEN, lw=1.3, radius=0.08, shadow=True)
text(s, sx, y4+0.13, sw, 0.28, [[("TALKER  →  output tokens", dict(size=12.5, color=GREEN, bold=True, font=BODY, spacing=0.6))]], align=PP_ALIGN.CENTER)
text(s, sx, y4+0.41, sw, 0.24, [[("autoregressive · sampling noise can't leak back into reasoning", dict(size=10, color=MUTE, font=BODY))]], align=PP_ALIGN.CENTER)

# ---------- cross-attention (read-only) dashed arrows ----------
arrow(s, ix+iw, 3.55, sx-0.05, 3.55, color=BLUE, w=1.4, dashed=True)
text(s, ix+iw+0.05, 3.28, sx-ix-iw, 0.24, [[("cross-attend · read-only", dict(size=9, color=BLUE, italic=True, font=BODY))]], align=PP_ALIGN.CENTER)
arrow(s, ix+iw, 5.9, sx-0.05, 5.9, color=BLUE, w=1.4, dashed=True)
text(s, ix+iw+0.05, 5.63, sx-ix-iw, 0.24, [[("raw tokens · exact quotes", dict(size=9, color=BLUE, italic=True, font=BODY))]], align=PP_ALIGN.CENTER)

# ---------- legend ----------
ly0 = 6.62
line(s, sx, ly0, sx+0.4, ly0, color=TEXT, w=2.0)
text(s, sx+0.5, ly0-0.12, 3.0, 0.24, [[("data flow · writes self-state", dict(size=9.5, color=MUTE, font=BODY))]])
line(s, sx+3.5, ly0, sx+3.9, ly0, color=BLUE, w=1.6, dash="dash")
text(s, sx+4.0, ly0-0.12, 3.0, 0.24, [[("read-only cross-attention", dict(size=9.5, color=MUTE, font=BODY))]])
pagenum(s, 4, "Proposed Solution · Pipeline")

# ============================================================
# SLIDE 5 — TRAINING PROCESS
# ============================================================
s = slide(BG)
eyebrow(s, "Training", color=AMBER)
title(s, "Two losses, staged so nothing collapses")
text(s, 0.7, 1.52, 12.0, 0.4, [[
    ("No component can be trained from a random init — each mechanism needs the previous one stable, so training walks a curriculum.",
     dict(size=13.5, color=MUTE, font=BODY))]], line=1.15)

# ----- left: the two-loss objective diagram -----
lx0 = 0.7
rect(s, lx0, 2.1, 5.75, 4.25, fill=CARD, lc=LINE, lw=1, radius=0.05, shadow=True)
text(s, lx0+0.32, 2.28, 5.1, 0.3, [[("ONE SHARED ENCODER, TWO LOSSES", dict(size=11.5, color=TEXT, bold=True, font=BODY, spacing=1.3))]])
# chunk -> shared encoder
rect(s, lx0+1.9, 2.72, 2.0, 0.46, fill=BG2, lc=LINE, lw=1, radius=0.14)
text(s, lx0+1.9, 2.72, 2.0, 0.46, [[("chunk cₜ", dict(size=11, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, lx0+2.9, 3.18, lx0+2.9, 3.42, color=MUTE, w=1.3)
rect(s, lx0+1.1, 3.46, 3.6, 0.5, fill=CARD2, lc=TEAL, lw=1.3, radius=0.12)
text(s, lx0+1.1, 3.46, 3.6, 0.5, [[("SHARED CHUNK ENCODER", dict(size=11, color=TEAL, bold=True, font=BODY, spacing=0.8))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# split to two branches
arrow(s, lx0+1.75, 3.96, lx0+1.2, 4.28, color=GREEN, w=1.3)
arrow(s, lx0+4.05, 3.96, lx0+4.6, 4.28, color=BLUE, w=1.3)
# branch A: reconstruction (anchor)
rect(s, lx0+0.3, 4.32, 2.55, 1.72, fill=BG2, lc=GREEN, lw=1.3, radius=0.08)
text(s, lx0+0.3, 4.46, 2.55, 0.5, [[("HRM loop", dict(size=11, color=TEXT, bold=True, font=BODY))],[("→ Talker → NLL", dict(size=11, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, line=1.0)
text(s, lx0+0.42, 5.12, 2.3, 0.5, [[("RECONSTRUCTION", dict(size=11.5, color=GREEN, bold=True, font=BODY, spacing=0.6))],[("decode the same chunk", dict(size=9.5, color=MUTE, font=BODY))]], align=PP_ALIGN.CENTER, line=1.05)
rect(s, lx0+0.42, 5.66, 2.3, 0.28, fill=None, lc=GREEN, lw=1, radius=0.5)
text(s, lx0+0.42, 5.66, 2.3, 0.28, [[("always-on anchor", dict(size=9, color=GREEN, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# branch B: SSL (secondary)
rect(s, lx0+2.95, 4.32, 2.55, 1.72, fill=BG2, lc=BLUE, lw=1.3, radius=0.08)
text(s, lx0+2.95, 4.46, 2.55, 0.5, [[("SSL proj head", dict(size=11, color=TEXT, bold=True, font=BODY))],[("→ cos vs EMA(cₜ₊₁)", dict(size=11, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, line=1.0)
text(s, lx0+3.07, 5.12, 2.3, 0.5, [[("SELF-SUPERVISED", dict(size=11.5, color=BLUE, bold=True, font=BODY, spacing=0.6))],[("predict the next chunk", dict(size=9.5, color=MUTE, font=BODY))]], align=PP_ALIGN.CENTER, line=1.05)
rect(s, lx0+3.07, 5.66, 2.3, 0.28, fill=None, lc=BLUE, lw=1, radius=0.5)
text(s, lx0+3.07, 5.66, 2.3, 0.28, [[("secondary · 0.1× · EMA 0.996", dict(size=8.5, color=BLUE, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# warm-up credit-assignment caption under the diagram
text(s, lx0+0.02, 6.46, 5.7, 0.3, [[
    ("Credit assignment: ", dict(size=10.5, color=TEXT, bold=True, font=BODY)),
    ("truncated BPTT warms up 2 → 5 steps (inner loop, then memory).", dict(size=10.5, color=MUTE, font=BODY))]])

# ----- right: the full curriculum, all stages listed -----
rx0 = 6.85; rw = 5.78
text(s, rx0, 2.16, rw, 0.32, [[("THE CURRICULUM — STAGES A → F", dict(size=13, color=AMBER, bold=True, font=BODY, spacing=1.3))]])
text(s, rx0, 2.48, rw, 0.28, [[("Nothing trains from a random init — each stage's stability unlocks the next.", dict(size=10.5, color=MUTE, italic=True, font=BODY))]])
stages5 = [
    (GREEN, "A", "Ground the Talker",     "latent → decodable text"),
    (GREEN, "B", "Inner HRM loop",        "fixed depth · detached mem"),
    (GREEN, "C", "Un-detach memory",      "grad reaches past thoughts"),
    (TEAL,  "D", "Self-supervised loss",  "JEPA cos vs EMA target"),
    (BLUE,  "E", "Adaptive depth (ACT)",  "halting · compute dial"),
    (AMBER, "F", "Chatbot fine-tune",     "two-lane · role tags"),
]
sy0 = 2.86; srh = 0.545; sgp = 0.07
for i,(col,code,name,desc) in enumerate(stages5):
    y = sy0 + i*(srh+sgp)
    rect(s, rx0, y, rw, srh, fill=CARD, lc=LINE, lw=1, radius=0.09, shadow=True)
    rect(s, rx0+0.16, y+0.11, 0.32, 0.32, fill=BG2, lc=col, lw=1.4, radius=0.5, shape=MSO_SHAPE.OVAL)
    text(s, rx0+0.16, y+0.11, 0.32, 0.32, [[(code, dict(size=13, color=col, bold=True, font=HEAD))]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx0+0.64, y, 2.55, srh, [[(name, dict(size=13, color=TEXT, bold=True, font=BODY))]], anchor=MSO_ANCHOR.MIDDLE)
    line(s, rx0+3.28, y+0.12, rx0+3.28, y+srh-0.12, color=LINE, w=1)
    text(s, rx0+3.44, y, rw-3.6, srh, [[(desc, dict(size=10.5, color=MUTE, font=BODY))]], anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 5, "Training")

# ============================================================
# SLIDE 6 — WHAT'S DIFFERENT (plain language)
# ============================================================
s = slide(BG)
eyebrow(s, "Proposed Solution", color=TEAL)
title(s, "What this does differently, in plain language")

rows = [
    (TEAL, "Deliberate, then decode",
     "Each thought is a bounded recurrent “think” — not one forward pass. Sampling noise from the Talker can never leak back into reasoning."),
    (BLUE, "Compute becomes a dial",
     "Adaptive halting spends more loop iterations on hard thoughts, fewer on filler. Parcae's constraint makes more looping predictably help, not add noise."),
    (GREEN, "Context at O(1), not O(n)",
     "Prior thoughts persist in a fixed-size gestalt memory that gradients still reach — long-range recall without a growing KV-cache."),
    (AMBER, "A self / input boundary",
     "Two lanes: user input can only be attended to; the reasoning loop alone writes self-state. The model can hold a belief distinct from the user's."),
]
gy = 1.72; rh = 1.16; rgap = 0.13
for i,(col,h,b) in enumerate(rows):
    y = gy + i*(rh+rgap)
    rect(s, 0.7, y, 12.08, rh, fill=CARD, lc=LINE, lw=1, radius=0.07, shadow=True)
    rect(s, 1.02, y+0.30, 0.56, 0.56, fill=BG2, lc=col, lw=1.5, radius=0.5)
    text(s, 1.02, y+0.30, 0.56, 0.56, [[("0%d"%(i+1), dict(size=15, color=col, bold=True, font=HEAD))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.95, y+0.20, 3.7, 0.8,
         [[(h, dict(size=16, color=TEXT, bold=True, font=HEAD))]],
         anchor=MSO_ANCHOR.MIDDLE, line=1.0)
    line(s, 5.9, y+0.22, 5.9, y+rh-0.22, color=LINE, w=1)
    text(s, 6.2, y+0.18, 6.4, 0.85,
         [[(b, dict(size=12.5, color=MUTE, font=BODY))]],
         anchor=MSO_ANCHOR.MIDDLE, line=1.18)
pagenum(s, 6, "Proposed Solution")

# ============================================================
# SLIDE 7 — PRELIMINARY FINDINGS
# ============================================================
s = slide(BG)
eyebrow(s, "Preliminary Findings", color=GREEN)
title(s, "A working end-to-end prototype — and one sharp result")
text(s, 0.7, 1.6, 12.1, 0.5, [[
    ("A from-scratch reference implementation runs the full A→F curriculum on real text and generates decodable output. Building it surfaced a real, non-obvious pathology.",
     dict(size=14, color=MUTE, font=BODY))]], line=1.2)

# left: the finding narrative
lx, ly, lw = 0.7, 2.4, 6.15
rect(s, lx, ly, lw, 4.05, fill=CARD, lc=LINE, lw=1, radius=0.05, shadow=True)
text(s, lx+0.36, ly+0.30, lw-0.72, 0.4,
     [[("The SSL-collapse pathology", dict(size=17, color=TEXT, bold=True, font=HEAD))]])
text(s, lx+0.36, ly+0.78, lw-0.72, 1.5, [
    [("Sharing one chunk encoder between the reconstruction loss and the self-supervised latent loss ",
      dict(size=12.5, color=MUTE, font=BODY)),
     ("collapses the latent", dict(size=12.5, color=CORAL, bold=True, font=BODY)),
     (": the encoder emits a near-constant vector for every chunk (cosine → 0.996), perfectly satisfying SSL while carrying no information.",
      dict(size=12.5, color=MUTE, font=BODY))],
], line=1.26)
text(s, lx+0.36, ly+2.16, lw-0.72, 0.3,
     [[("Why it's dangerous", dict(size=11, color=GREEN, bold=True, font=BODY, spacing=1.4))]])
for j,(t) in enumerate([
    "Silent — SSL→0 looks exactly like success",
    "Propagating — a shared encoder spreads the damage",
    "Absorbing — the EMA target is a stable dead end"]):
    yy = ly+2.52+j*0.44
    dot(s, lx+0.4, yy+0.05, 0.11, GREEN)
    text(s, lx+0.66, yy-0.04, lw-0.9, 0.32, [[(t, dict(size=11.5, color=TEXT, font=BODY))]],
         line=1.0, anchor=MSO_ANCHOR.MIDDLE)

# right: before/after stat cards
rx = 7.1
text(s, rx, 2.28, 5.7, 0.3, [[("BEFORE → AFTER THE FIX", dict(size=11.5, color=MUTE, bold=True, font=BODY, spacing=1.6))]])
stats = [
    ("Validation loss, D/E", "8.26", "~7.8", CORAL, GREEN, "regressed", "held flat"),
    ("Latent representation", "cos 0.996", "std ~0.25", CORAL, GREEN, "collapsed", "healthy"),
    ("Reconstruction ppl", "~3360", "~3136", MUTE, GREEN, "broken run", "fixed run"),
]
sy = 2.66; sh_ = 1.16; sg = 0.12
for lab, before, after, cb, ca, tb, ta in stats:
    rect(s, rx, sy, 5.68, sh_, fill=CARD, lc=LINE, lw=1, radius=0.07, shadow=True)
    text(s, rx+0.34, sy+0.18, 5.0, 0.3, [[(lab, dict(size=12, color=TEXT, bold=True, font=BODY))]])
    text(s, rx+0.34, sy+0.5, 2.55, 0.5,
         [[(before, dict(size=19, color=cb, bold=True, font=HEAD))]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx+0.34, sy+0.92, 2.55, 0.2, [[(tb, dict(size=9, color=FAINT, font=BODY))]])
    # arrow
    text(s, rx+2.9, sy+0.5, 0.5, 0.5, [[("→", dict(size=18, color=MUTE, font=BODY))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx+3.45, sy+0.5, 2.1, 0.5,
         [[(after, dict(size=19, color=ca, bold=True, font=HEAD))]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx+3.45, sy+0.92, 2.1, 0.2, [[(ta, dict(size=9, color=GREEN, font=BODY))]])
    sy += sh_ + sg
text(s, rx, sy+0.02, 5.68, 0.4, [[
    ("The fix: reconstruction is the always-on anchor; SSL demoted to its own head + variance floor + a collapse monitor.",
     dict(size=10.5, color=FAINT, italic=True, font=BODY))]], line=1.15)
pagenum(s, 7, "Preliminary Findings")

# ============================================================
# SLIDE 8 — BASELINE COMPARISON
# ============================================================
s = slide(BG)
eyebrow(s, "Baseline Comparison", color=BLUE)
title(s, "At matched compute, competitive with a standard GPT")
text(s, 0.7, 1.5, 12.1, 0.4, [[
    ("Memorizing one Wikipedia page — same gpt2 tokenizer, optimizer, schedule & step budget; ", dict(size=13, color=MUTE, font=BODY)),
    ("only the architecture differs.", dict(size=13, color=TEXT, bold=True, font=BODY))]], line=1.15)

# --- the comparison chart (runs/comparison.png, 1425x870) ---
img_h = 4.32
img_w = img_h * (1425/870.0)      # ≈ 7.08
imgx, imgy = 0.62, 2.02
rect(s, imgx-0.08, imgy-0.08, img_w+0.16, img_h+0.16, fill="F4F6FA", lc=LINE, lw=1, radius=0.02, shadow=True)
s.shapes.add_picture("runs/comparison.png", Inches(imgx), Inches(imgy), Inches(img_w), Inches(img_h))

# --- interpretation panel (right) ---
px = imgx + img_w + 0.32          # ≈ 8.31
pw = 12.63 - px                   # ≈ 4.32
text(s, px, 2.02, pw, 0.3, [[("WHAT IT SHOWS", dict(size=12, color=TEXT, bold=True, font=BODY, spacing=1.5))]])
takes = [
    (GREEN, "Both memorize the page",
     "Latent grounded-only reconstruction → ppl ≈ 1.0, matching the 44.7M same-params GPT (≈ 1.1)."),
    (CORAL, "Matched-width GPT can't",
     "The 14.1M same-compute GPT (d=192) plateaus at ≈ 484 — same width, no bottleneck to lean on."),
    (AMBER, "The SSL cost is visible",
     "Under the full A→E curriculum the latent model sits at 936: self-supervision trades page-fit for generalization."),
]
ty = 2.46; trh = 1.16
for i,(col,h,b) in enumerate(takes):
    y = ty + i*trh
    dot(s, px+0.02, y+0.06, 0.16, col)
    text(s, px+0.3, y-0.04, pw-0.3, 0.32, [[(h, dict(size=13, color=TEXT, bold=True, font=HEAD))]], line=1.0)
    text(s, px+0.3, y+0.3, pw-0.3, 0.8, [[(b, dict(size=11, color=MUTE, font=BODY))]], line=1.16)
    if i < len(takes)-1:
        line(s, px+0.3, y+trh-0.14, px+pw, y+trh-0.14, color=LINE, w=1)

# honest-read strip under the chart (one line)
by = imgy + img_h + 0.2
text(s, imgx, by, img_w+0.2, 0.3, [[
    ("Chance ≈ 50,262.  ", dict(size=10, color=FAINT, bold=True, font=BODY)),
    ("One-page memorization isn't the goal — the payoff shows only at scale.",
     dict(size=10, color=FAINT, italic=True, font=BODY))]], line=1.1)
pagenum(s, 8, "Baseline Comparison")

# ============================================================
# SLIDE 9 — SUCCESS METRICS
# ============================================================
s = slide(BG)
eyebrow(s, "Success Metrics", color=TEAL)
title(s, "Three signals decide whether it works")

metrics = [
    (TEAL, "1", "Decodability",
     "Reconstruction\nperplexity",
     "Can the Talker faithfully rebuild a chunk from its latent? The anchor metric — it cannot be gamed by a degenerate latent."),
    (BLUE, "2", "Representation health",
     "latent_std\ncollapse monitor",
     "Per-dimension latent variance, logged every eval. Must stay above the floor — the tripwire that would catch any SSL collapse early."),
    (AMBER, "3", "Compute scaling",
     "Loss vs. loop\niterations",
     "Does more test-time looping help along Parcae's predictable, saturating curve? Confirms depth is a real, monotone compute dial."),
]
cw, gap = 3.86, 0.35
cx = 0.7; cy = 1.95; ch = 4.35
for col, num, name, metric, desc in metrics:
    rect(s, cx, cy, cw, ch, fill=CARD, lc=LINE, lw=1, radius=0.06, shadow=True)
    rect(s, cx+0.36, cy+0.4, 0.7, 0.7, fill=BG2, lc=col, lw=1.6, radius=0.5)
    text(s, cx+0.36, cy+0.4, 0.7, 0.7, [[(num, dict(size=22, color=col, bold=True, font=HEAD))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.36, cy+1.32, cw-0.72, 0.4,
         [[(name.upper(), dict(size=11.5, color=col, bold=True, font=BODY, spacing=1.4))]])
    mlines = metric.split("\n")
    text(s, cx+0.36, cy+1.72, cw-0.72, 0.9,
         [[(mlines[0], dict(size=19, color=TEXT, bold=True, font=HEAD))],
          [(mlines[1], dict(size=19, color=TEXT, bold=True, font=HEAD))]], line=1.02)
    line(s, cx+0.36, cy+2.72, cx+cw-0.36, cy+2.72, color=LINE, w=1)
    text(s, cx+0.36, cy+2.9, cw-0.72, 1.3,
         [[(desc, dict(size=12.5, color=MUTE, font=BODY))]], line=1.24)
    cx += cw + gap
pagenum(s, 9, "Success Metrics")

# ============================================================
# SLIDE 10 — PLAN: TIMELINE + RISKS
# ============================================================
s = slide(BG)
eyebrow(s, "Plan", color=BLUE)
title(s, "Curriculum milestones, then risks")
text(s, 0.7, 1.58, 12, 0.4, [[
    ("The architecture can't be trained from a random init — each mechanism needs the previous one stable. That staging is the roadmap.",
     dict(size=13.5, color=MUTE, font=BODY))]], line=1.2)

stages = [
    ("A", "Ground the Talker", "latent → text", GREEN, True),
    ("B", "Inner HRM loop", "fixed depth", GREEN, True),
    ("C", "Un-detach memory", "cross-thought grad", GREEN, True),
    ("D", "Add SSL loss", "collapse fix here", TEAL, True),
    ("E", "Adaptive depth", "ACT halting", BLUE, False),
    ("F", "Chatbot fine-tune", "two-lane / roles", MUTE, False),
]
# timeline
tx, ty = 0.7, 2.35
seg = 12.08/len(stages)
line(s, tx+0.2, ty+0.55, tx+12.08-0.2, ty+0.55, color=LINE, w=2)
for i,(code,name,sub,col,done) in enumerate(stages):
    cxp = tx + i*seg + seg/2
    dd = 0.5
    rect(s, cxp-dd/2, ty+0.3, dd, dd, fill=(col if done else BG2), lc=col, lw=2, radius=0.5, shape=MSO_SHAPE.OVAL)
    text(s, cxp-dd/2, ty+0.3, dd, dd, [[(code, dict(size=15, color=(BG if done else col), bold=True, font=HEAD))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxp-seg/2+0.1, ty+0.98, seg-0.2, 0.4,
         [[(name, dict(size=12, color=TEXT, bold=True, font=BODY))]], align=PP_ALIGN.CENTER, line=1.0)
    text(s, cxp-seg/2+0.1, ty+1.36, seg-0.2, 0.3,
         [[(sub, dict(size=10, color=col, italic=True, font=BODY))]], align=PP_ALIGN.CENTER)
text(s, tx, ty-0.28, 12, 0.3, [[
    ("Stages A–D verified on real text (≈1.5M tokens) · E–F wired but not yet scaled → next: a GPU shakedown before the first large run",
     dict(size=11, color=FAINT, italic=True, font=BODY))]], align=PP_ALIGN.CENTER)

# risks
ry = 4.55
text(s, 0.7, ry-0.06, 6, 0.3, [[("TOP RISKS & MITIGATIONS", dict(size=12, color=AMBER, bold=True, font=BODY, spacing=1.6))]])
risks = [
    (CORAL, "SSL collapse returns at scale",
     "Hyperparameters were tuned on a 1.5M-token smoke run.",
     "Reconstruction stays the always-on anchor; separate SSL head + variance floor + latent_std monitor catch any recurrence early."),
    (AMBER, "No large run has been done",
     "Everything is verified at smoke/offline scale; a global LR cosine starves late stages.",
     "Per-stage warmup→cosine schedule; short GPU shakedown to confirm throughput and health before committing to a long run."),
]
rcw = 5.86; rgpx = 0.35; rhh = 1.9
for i,(col,h,risk,mit) in enumerate(risks):
    x = 0.7 + i*(rcw+rgpx)
    rect(s, x, ry+0.32, rcw, rhh, fill=CARD, lc=LINE, lw=1, radius=0.06, shadow=True)
    dot(s, x+0.34, ry+0.66, 0.2, col)
    text(s, x+0.7, ry+0.5, rcw-1.0, 0.4, [[(h, dict(size=15, color=TEXT, bold=True, font=HEAD))]], line=1.0)
    text(s, x+0.34, ry+1.02, rcw-0.68, 0.5,
         [[("Risk  ", dict(size=10.5, color=col, bold=True, font=BODY)),
           (risk, dict(size=11.5, color=MUTE, font=BODY))]], line=1.15)
    text(s, x+0.34, ry+1.5, rcw-0.68, 0.5,
         [[("Mitigation  ", dict(size=10.5, color=GREEN, bold=True, font=BODY)),
           (mit, dict(size=11.5, color=TEXT, font=BODY))]], line=1.15)
pagenum(s, 10, "Plan")

# ============================================================
# SLIDE 11 — CLOSING
# ============================================================
s = slide(BG2)
# rings motif upper-right (empty zone; mirrors the title slide, clear of text)
_rc_x, _rc_y = 12.9, 2.55
for i,(d,a) in enumerate([(3.4,"16000"),(2.5,"26000"),(1.5,"42000")]):
    ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(_rc_x-d/2), Inches(_rc_y-d/2), Inches(d),Inches(d))
    ov.fill.background(); ov.line.color.rgb=C(TEAL); ov.line.width=Pt(1.1)
    ln = ov.line._get_or_add_ln()
    af = ln.makeelement(qn('a:solidFill'),{}); ce=af.makeelement(qn('a:srgbClr'),{'val':TEAL})
    al=ce.makeelement(qn('a:alpha'),{'val':a}); ce.append(al); af.append(ce)
    for ex in ln.findall(qn('a:solidFill')): ln.remove(ex)
    ln.insert(0,af); ov.shadow.inherit=False

eyebrow(s, "Summary", x=0.9, y=2.1, color=TEAL)
text(s, 0.85, 2.6, 11.6, 2.0, [
    [("Think first.", dict(size=44, color=TEXT, bold=True, font=HEAD))],
    [("Speak second.", dict(size=44, color=TEAL, bold=True, font=HEAD))],
], line=1.05)
text(s, 0.9, 4.55, 10.8, 1.0, [[
    ("One architecture that separates reasoning from expression, turns deliberation into a tunable compute dial, and keeps a self / input boundary — with a working prototype and a de-risked path to scale.",
     dict(size=16, color=MUTE, font=BODY))]], line=1.32)
line(s, 0.9, 6.15, 12.43, 6.15, color=LINE, w=1)
text(s, 0.9, 6.35, 11.5, 0.4, [[
    ("Hierarchical Latent Reasoning Architecture", dict(size=13, color=TEXT, bold=True, font=BODY)),
    ("   ·   Project Proposal", dict(size=13, color=FAINT, font=BODY))]])

prs.save("/Users/xiajin/djin/code/ai_projects/ucsc/Hierarchical_Latent_Reasoning_Architecture.pptx")
print("saved")
