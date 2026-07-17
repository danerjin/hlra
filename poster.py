#!/usr/bin/env python
"""ISEF tri-fold poster — 48 x 36 in canvas, content kept OFF the hinges
(x=12, x=36). Columns sit inside the three panels; the 4x4 letter tiling maps
1 col (left) + 2 cols (center) + 1 col (right) x 4 rows, with the hinge gaps and
edge margins left blank. Pale (light) card fills, colored borders, arrows on top."""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INK="10151F"; MUTE="53606E"; FAINT="93A0AE"
TEAL="0E7C6B"; BLUE="1D4ED8"; AMBER="B4560A"; NAVY="0B3B6F"
SOFT="C2CAD6"
PALE={TEAL:"EFF7F5",BLUE:"EEF3FD",AMBER:"FBF5EA",NAVY:"EFF3F8",SOFT:"F7F9FB"}
SERIF="Georgia"; SANS="Calibri"
def C(h): return RGBColor.from_string(h)

prs=Presentation(); prs.slide_width=Inches(48); prs.slide_height=Inches(36)
s=prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb=C("FFFFFF")

def rect(x,y,w,h,fill=None,lc=None,lw=1.0,radius=0.05,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if lc is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(lc); sp.line.width=Pt(lw)
    if radius is not None and shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    sp.shadow.inherit=False; return sp

def T(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line=1.08,sa=5,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line
        if sa is not None: p.space_after=Pt(sa)
        for tx,o in para:
            r=p.add_run(); r.text=tx
            r.font.size=Pt(o.get("size",22)); r.font.bold=o.get("bold",False)
            r.font.italic=o.get("italic",False); r.font.name=o.get("font",SANS)
            r.font.color.rgb=C(o.get("color",INK))
            sp=o.get("spacing")
            if sp is not None: r._r.get_or_add_rPr().set('spc',str(int(sp*100)))
    return tb

def line(x1,y1,x2,y2,color=INK,w=2.4,dash=None):
    ln=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=C(color); ln.line.width=Pt(w)
    if dash: d=ln.line._get_or_add_ln(); d.append(d.makeelement(qn('a:prstDash'),{'val':dash}))
    ln.shadow.inherit=False; return ln

def arrow(x1,y1,x2,y2,color=INK,w=3.0,dash=None,head=0.42):
    line(x1,y1,x2,y2,color=color,w=w,dash=dash)
    dx,dy=x2-x1,y2-y1
    rot=(90 if dx>=0 else 270) if abs(dx)>=abs(dy) else (180 if dy>=0 else 0)
    tri=s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,Inches(x2-head/2),Inches(y2-head/2),Inches(head),Inches(head))
    tri.rotation=rot; tri.fill.solid(); tri.fill.fore_color.rgb=C(color)
    tri.line.fill.background(); tri.shadow.inherit=False

def dot(x,y,d,color): return rect(x,y,d,d,fill=color,radius=None,shape=MSO_SHAPE.OVAL)

def card(x,y,w,h,accent,title,body,tsize=25,bsize=21,lw=2.0,pad=0.4,center=False,tcolor=None):
    rect(x,y,w,h,fill=PALE.get(accent,"FFFFFF"),lc=accent,lw=lw,radius=0.05)
    al=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    runs=[]
    if title: runs.append([(title,dict(size=tsize,bold=True,color=tcolor or accent))])
    if isinstance(body,str): body=[[(body,dict(size=bsize,color=INK))]]
    runs+=body
    T(x+pad,y+0.18,w-2*pad,h-0.36,runs,align=al,anchor=MSO_ANCHOR.MIDDLE,line=1.12,sa=7)

def section_header(x,y,w,label,size=42):
    dot(x,y+0.10,0.40,NAVY)
    T(x+0.66,y-0.10,w-0.66,0.95,[[(label.upper(),dict(size=size,bold=True,color=NAVY,font=SERIF,spacing=0.3))]])
    line(x,y+1.18,x+w,y+1.18,color=SOFT,w=2.0)

def measure(sec):
    return 1.5+sum(it["h"] for it in sec["items"])+0.34*len(sec["items"])
def draw_column(x,w,y0,y1,secs,inner=0.34):
    hs=[measure(sc) for sc in secs]; gap=(y1-y0-sum(hs))/max(1,len(secs)-1); y=y0
    for sc,H in zip(secs,hs):
        section_header(x,y,w,sc["title"]); yy=y+1.5
        for it in sc["items"]: it["draw"](x,yy,w,it["h"]); yy+=it["h"]+inner
        y+=H+gap

def it_lead(runs,h=1.3): return {"h":h,"draw":lambda x,y,w,hh:T(x,y,w,hh,runs,line=1.14,sa=0)}
def it_card(accent,title,body,h,**kw): return {"h":h,"draw":lambda x,y,w,hh:card(x,y,w,hh,accent,title,body,**kw)}
def it_numcard(n,accent,title,body,h):
    def d(x,y,w,hh):
        rect(x,y,w,hh,fill=PALE[SOFT],lc=SOFT,lw=1.7,radius=0.05)
        rect(x+0.38,y+(hh-1.0)/2,1.0,1.0,fill="FFFFFF",lc=accent,lw=2.4,radius=0.16)
        T(x+0.38,y+(hh-1.0)/2,1.0,1.0,[[(n,dict(size=32,bold=True,color=accent,font=SERIF))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
        T(x+1.68,y+0.2,w-2.0,hh-0.4,[[(title,dict(size=25,bold=True,color=INK))],[(body,dict(size=19.5,color=INK))]],anchor=MSO_ANCHOR.MIDDLE,line=1.12,sa=7)
    return {"h":h,"draw":d}
def it_bullets(rows,h,bs=21):
    def d(x,y,w,hh):
        step=hh/len(rows)
        for i,(col,txt) in enumerate(rows):
            yy=y+i*step; dot(x+0.12,yy+step/2-0.11,0.24,col)
            T(x+0.6,yy,w-0.7,step,[[(txt,dict(size=bs,color=INK))]],anchor=MSO_ANCHOR.MIDDLE,line=1.06,sa=0)
    return {"h":h,"draw":d}

# ---- figure slots -----------------------------------------------------------
# `python poster_figs.py runs/scaled` renders these PNGs at exactly the slot size
# (4.8 x 2.55 in). If a PNG is present it is dropped in; if not, the slot keeps its
# "data coming soon" placeholder, so the poster always builds.
FIGS=os.path.join(os.path.dirname(os.path.abspath(__file__)),"poster_figs")
def figslot(x,y,w,hh,png,title):
    path=os.path.join(FIGS,png)
    if os.path.exists(path):
        s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(hh)); return
    rect(x,y,w,hh,fill="FFFFFF",lc=SOFT,lw=1.6,radius=0.03)
    line(x+0.55,y+hh-0.55,x+w-0.45,y+hh-0.55,color=SOFT,w=1.2)
    line(x+0.65,y+0.5,x+0.65,y+hh-0.55,color=SOFT,w=1.2)
    T(x,y+0.4,w,0.4,[[(title,dict(size=17,bold=True,color=MUTE))]],align=PP_ALIGN.CENTER,sa=0)
    T(x,y+hh-0.5,w,0.35,[[("data coming soon",dict(size=13.5,italic=True,color=FAINT))]],align=PP_ALIGN.CENTER,sa=0)
def it_figs(h=2.55):
    def d(x,y,w,hh):
        half=(w-0.4)/2                       # 4.8in -- must match poster_figs.FIG_W
        figslot(x,y,half,hh,"fig_loss.png","Loss vs. steps")
        figslot(x+half+0.4,y,half,hh,"fig_arc.png","ARC-C vs. FLOPs")
    return {"h":h,"draw":d}

# ---- panel geometry: hinges at x=12 & x=36; keep content off them ----
MTOP=1.3; MBOT=34.6
LX,LW = 1.0, 10.0            # left panel  (tile band 0.5–11.5)
RX,RW = 37.0, 10.0          # right panel (tile band 36.5–47.5)
CX,CW = 13.3, 21.4          # center panel (tile bands 13–24, 24–35)

# =====================================================================
# LEFT COLUMN
# =====================================================================
draw_column(LX,LW,MTOP,MBOT,[
 {"title":"The problem","items":[
   it_lead([[("Chain-of-thought makes a model externalize every reasoning step as tokens. Three concrete costs:",dict(size=22,color=INK))]],1.35),
   it_numcard("1",TEAL,"Thought = word-piece","Each step must be a valid token, so the unit of reasoning is a sub-word — and error compounds at that fine grain.",2.6),
   it_numcard("2",BLUE,"Compute is fixed","≈ one forward pass per token. No way to spend more thinking on a hard step without emitting more text.",2.6),
   it_numcard("3",AMBER,"Reasoning ⊕ phrasing","No representational place to hold “what I concluded” apart from “how I’d say it.”",2.6),
 ]},
 {"title":"The goal","items":[
   it_lead([[("Make the unit of reasoning a whole semantic chunk — a latent vector, not a token span:",dict(size=22,color=INK))]],1.2),
   it_bullets([(BLUE,"A bounded fast/slow recurrent loop refines and predicts these thoughts."),
               (BLUE,"It reads a differentiable memory of earlier thoughts — credit for a later conclusion flows back into earlier ones."),
               (BLUE,"Adaptive depth: iterate more on a hard thought without producing any extra output.")],3.9),
   it_card(AMBER,"THE OPEN QUESTION",
     [[("Can a model reason coherently and forward in a compact latent space — and translate to language only at the end — instead of thinking in the token stream?",dict(size=21,italic=True,color=INK))]],2.9,tsize=19,lw=2.4),
 ]},
 {"title":"Hypothesis","items":[
   it_lead([[("Latent-space reasoning is trainable ",dict(size=22,color=INK)),("if and only if",dict(size=22,bold=True,color=INK)),(" two failure modes are held off at once:",dict(size=22,color=INK))]],1.35),
   it_card(TEAL,"Representational collapse","a self-predictive latent objective has a trivial fixed point — predict a constant from a constant.",2.0,tsize=23),
   it_card(BLUE,"Recurrent instability","an unbounded reasoning loop drifts or explodes as it iterates deeper.",2.0,tsize=23),
   it_card(TEAL,"THE CLAIM",
     [[("A reconstruction anchor + a norm-bounded loop + a differentiable thought-memory ",dict(size=21,bold=True,color=INK)),
       ("jointly",dict(size=21,bold=True,italic=True,color=TEAL)),(" make latent reasoning stable enough to train.",dict(size=21,bold=True,color=INK))]],2.5,tsize=21,lw=2.4),
 ]},
])

# =====================================================================
# RIGHT COLUMN
# =====================================================================
draw_column(RX,RW,MTOP,MBOT,[
 {"title":"Why it’s new","items":[
   it_lead([[("Not any one of the four ideas — the claim that they are ",dict(size=22,color=INK)),("mutually load-bearing:",dict(size=22,bold=True,color=INK))]],1.3),
   it_bullets([(TEAL,"JEPA latent prediction alone → collapses."),
               (BLUE,"An HRM-style loop alone → nothing anchors its latents to language."),
               (AMBER,"A differentiable thought-memory alone → unstable to train through.")],3.6,bs=20),
   it_card(NAVY,"CONCRETELY NOVEL",
     [[("(1) ",dict(size=20,bold=True,color=TEAL)),("a disjoint two-objective split over a shared encoder as an anti-collapse strategy;",dict(size=20,color=INK))],
      [("(2) ",dict(size=20,bold=True,color=BLUE)),("reasoning at whole-chunk granularity with un-detached cross-thought credit;",dict(size=20,color=INK))],
      [("(3) ",dict(size=20,bold=True,color=AMBER)),("a reasoner that varies its own depth without emitting a token.",dict(size=20,color=INK))]],4.0,tsize=23),
 ]},
 {"title":"Results — in progress","items":[
   it_lead([[("Big training run underway. Planned reporting:",dict(size=21,color=INK))]],0.7),
   it_bullets([(NAVY,"Training curves — val loss + latent_std collapse monitor across the staged curriculum."),
               (NAVY,"ARC-C accuracy on a FLOPs / parameter scale, vs. other LLMs.")],2.7),
   it_figs(2.55),
   it_card(TEAL,"EARLY EVIDENCE (smoke scale)",
     [[("Full curriculum trains end-to-end; no collapse — val loss keeps ",dict(size=19.5,color=INK)),
       ("falling",dict(size=19.5,bold=True,color=INK)),
       (" when prediction turns on, where a collapse would show as a ",dict(size=19.5,color=INK)),
       ("rise",dict(size=19.5,bold=True,color=INK)),
       ("; grounded reconstruction matches a same-parameter GPT on a memorization probe.",dict(size=19.5,color=INK))]],3.0,tsize=21),
 ]},
 {"title":"Further research","items":[
   it_card(AMBER,"Make “think harder” learnable","an RL / supervised halt gate (today ACT degenerates to minimum depth).",2.2,tsize=22,bsize=19.5),
   it_card(BLUE,"Two-lane input / self memory","for dialogue — anti-sycophancy (“user said X” vs. “I concluded X”).",2.2,tsize=22,bsize=19.5),
   it_card(TEAL,"A successor architecture","(spec in progress) built directly on latent-space reasoning.",2.2,tsize=22,bsize=19.5),
 ]},
])

# =====================================================================
# CENTER
# =====================================================================
T(CX,1.35,CW,2.2,[
 [("Hierarchical Latent ",dict(size=54,bold=True,color=NAVY,font=SERIF)),("Reasoning",dict(size=54,bold=True,color=TEAL,font=SERIF))],
 [("Architecture",dict(size=54,bold=True,color=NAVY,font=SERIF))]],align=PP_ALIGN.CENTER,line=0.98,sa=0)
T(CX+0.9,4.0,CW-1.8,1.5,[[
 ("A language model whose reasoning happens in a recurrent loop over chunk-level latent vectors — one vector per sentence-like ",dict(size=23,color=INK)),
 ("“thought”",dict(size=23,bold=True,color=TEAL)),
 (" — with a separate decoder turning a thought into words only after the reasoning is done, so the reasoner never generates a token to think.",dict(size=23,color=INK))]],
 align=PP_ALIGN.CENTER,line=1.14,sa=0)
line(CX+2.3,5.75,CX+CW-2.3,5.75,color=NAVY,w=2.4)

Dx,Dw = CX, 12.4
Mx,Mw = CX+12.9, 8.5
section_header(Dx,6.1,Dw,"The architecture",size=34)
section_header(Mx,6.1,Mw,"Method",size=34)

# ---------- ARCHITECTURE: single top->bottom spine (boxes, then arrows) ----------
Sc=Dx+Dw/2
def fb(x,y,w,h,title,sub=None,lc="3A4453",lw=2.2,tc=INK,ts=21,ss=15,rad=0.08,fill="FFFFFF"):
    rect(x,y,w,h,fill=fill,lc=lc,lw=lw,radius=rad)
    if sub is None:
        T(x,y,w,h,[[(title,dict(size=ts,bold=True,color=tc))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
    else:
        T(x,y,w,h,[[(title,dict(size=ts,bold=True,color=tc))],[(sub,dict(size=ss,color=MUTE))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line=1.04,sa=2)

Btop=7.65; cur=[Btop]; ycen=[]
def sbox(w,h,*a,**k):
    y=cur[0]; fb(Sc-w/2,y,w,h,*a,**k); cur[0]=y+h; return y
def sgap(g=0.95,color=INK,dash=None):
    y=cur[0]; ycen.append((y,y+g,color,dash)); cur[0]=y+g

sbox(8.4,1.1,"Input text","“Every effort moves you. It compounds.”",ts=19,ss=14); sgap()
sbox(7.0,1.05,"SaT-Capped chunker","sentence / clause chunks",ts=19,ss=14); sgap()
cy=cur[0]
for lab,dx in (("c₁",-2.4),("c₂",0),("c₃",2.4)): fb(Sc+dx-1.05,cy,2.1,0.82,lab,ts=16,rad=0.14)
cur[0]=cy+0.82; sgap()
sbox(8.6,1.15,"CHUNK ENCODER","bidirectional transformer + masked mean-pool",ts=21,ss=14); sgap()
zy=cur[0]
rect(Sc-2.25,zy,4.5,0.95,fill="FFFFFF",lc=NAVY,lw=3.0,radius=0.5)
T(Sc-2.25,zy,4.5,0.95,[[("thought  zₜ",dict(size=21,bold=True,color=NAVY))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sa=0)
cur[0]=zy+0.95
rect(Dx+0.2,zy-0.72,2.95,2.35,fill=PALE[TEAL],lc=TEAL,lw=2.0,radius=0.06)
T(Dx+0.3,zy-0.6,2.75,2.1,[[("RECONSTRUCTION",dict(size=15,bold=True,color=TEAL))],
    [("the Talker reproduces this chunk — the anti-collapse ",dict(size=14,color=INK)),("anchor",dict(size=14,bold=True,color=INK)),(" (a constant can’t).",dict(size=14,color=INK))]],anchor=MSO_ANCHOR.MIDDLE,line=1.06,sa=3)
T(Sc+2.5,zy-0.05,2.9,1.1,[[("d_latent = mult × d_model",dict(size=14,bold=True,color=AMBER))],[("a thought is wider than a token",dict(size=13,color=INK))]],line=1.04,sa=0)
sgap()
Ly=cur[0]; Lw=6.9; Lh=4.4
rect(Sc-Lw/2,Ly,Lw,Lh,fill=PALE[BLUE],lc=BLUE,lw=2.8,radius=0.04)
T(Sc-Lw/2,Ly+0.16,Lw,0.5,[[("INNER HRM LOOP  ·  the reasoner",dict(size=19,bold=True,color=BLUE))]],align=PP_ALIGN.CENTER,sa=0)
mx=Sc-1.9; my=Ly+0.8
for k in range(3): fb(mx,my,0.7,0.6,"L",ts=16,rad=0.16); mx+=0.83
fb(mx+0.1,my,0.84,0.6,"H",lc=NAVY,tc=NAVY,ts=16,rad=0.16)
T(Sc-Lw/2,my+0.64,Lw,0.4,[[("fast L ×3  →  slow H ×1",dict(size=14,color=MUTE))]],align=PP_ALIGN.CENTER,sa=0)
for txt,col,dyy in [("hard-norm shell → bounded at any depth",TEAL,2.15),("decay gate a∈(0,1) → converges, not drifts",BLUE,2.66),("ACT halting → variable depth per thought",AMBER,3.17)]:
    dot(Sc-Lw/2+0.5,Ly+dyy+0.05,0.15,col); T(Sc-Lw/2+0.78,Ly+dyy-0.06,Lw-1.05,0.5,[[(txt,dict(size=14,color=INK))]],line=0.98,sa=0)
cur[0]=Ly+Lh
Mmx=Sc+Lw/2+0.28
rect(Mmx,Ly+0.5,Dx+Dw-0.35-Mmx,2.6,fill=PALE[BLUE],lc=BLUE,lw=2.0,radius=0.08)
T(Mmx,Ly+0.5,Dx+Dw-0.35-Mmx,2.6,[[("GESTALT",dict(size=15,bold=True,color=BLUE))],[("MEMORY",dict(size=15,bold=True,color=BLUE))],
    [("FIFO ·",dict(size=12.5,color=INK))],[("role-tagged",dict(size=12.5,color=INK))],[("un-detached",dict(size=12.5,bold=True,color=BLUE))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line=1.04,sa=1)
sgap(color=BLUE)
sbox(7.6,1.3,"pred_head  →  next thought  ẑₜ₊₁","≈ slow EMA target (JEPA scaled-cosine)",lc=BLUE,tc=BLUE,ts=18,ss=13); sgap(color=BLUE)
sbox(6.0,1.1,"TALKER  ·  decode","turns the finished thought into words",lc=AMBER,tc=AMBER,ts=19,ss=13); sgap(color=AMBER)
wy=cur[0]
rect(Sc-4.5,wy,9.0,1.3,fill=PALE[AMBER],lc=AMBER,lw=3.0,radius=0.06)
T(Sc-4.5,wy,9.0,1.3,[[("WORDS OUT",dict(size=21,bold=True,color=AMBER))],[("the only place tokens are ever produced",dict(size=14.5,color=INK))]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line=1.04,sa=2)
# arrows ON TOP
for y0,y1,col,dash in ycen: arrow(Sc,y0-0.02,Sc,y1+0.02,color=col,dash=dash,head=0.44)
arrow(Sc-2.25,zy+0.475,Dx+3.15,zy+0.475,color=TEAL,w=2.6,head=0.32)
arrow(Mmx,Ly+1.4,Sc+Lw/2-0.02,Ly+1.4,color=BLUE,w=2.4,head=0.30)
arrow(Sc+Lw/2,Ly+2.45,Mmx+0.02,Ly+2.45,color=BLUE,w=2.4,head=0.30)

# ---------- METHOD ----------
def mcard(y,h,accent,title,body,tsize=23): card(Mx,y,Mw,h,accent,title,body,tsize=tsize,bsize=20,lw=2.2,pad=0.36)
T(Mx,7.55,Mw,1.55,[[("Text → chunks → a bidirectional encoder gives one latent ",dict(size=21,color=INK)),
    ("thought",dict(size=21,bold=True,color=TEAL)),(" per chunk. Two objectives share the encoder but touch the rest ",dict(size=21,color=INK)),
    ("disjointly:",dict(size=21,bold=True,color=INK))]],line=1.12,sa=0)
mcard(9.35,3.7,TEAL,"RECONSTRUCTION — the anchor",
     [[("encoder → Talker decodes the chunk. A constant can’t reconstruct varied chunks ⇒ always-on ",dict(size=20,color=INK)),
       ("anti-collapse anchor",dict(size=20,bold=True,color=INK)),(" + a clean latent→text decoder.",dict(size=20,color=INK))]])
mcard(13.4,4.75,BLUE,"PREDICTION — the reasoning",
     [[("a bounded fast/slow loop runs over the thoughts, reads/writes a FIFO gestalt memory, and predicts the next thought vs a slow EMA target. Memory ",dict(size=20,color=INK)),
       ("un-detached",dict(size=20,bold=True,color=INK)),(" ⇒ cross-thought credit; decay gate + norm-shell ⇒ ",dict(size=20,color=INK)),
       ("stable at any depth",dict(size=20,bold=True,color=INK)),("; ACT head ⇒ variable depth.",dict(size=20,color=INK))]])
mcard(18.45,3.55,AMBER,"THE LOAD-BEARING MOVE",
     [[("The loop lives in prediction ",dict(size=20,color=INK)),("only",dict(size=20,bold=True,italic=True,color=AMBER)),
       (" — never reconstruction — so a thought isn’t torn between “decode me now” and “predict what’s next.”",dict(size=20,color=INK))]],tsize=22)
mcard(22.3,3.3,NAVY,"CURRICULUM — staged",
     [[("autoencoder → loop → un-detach memory → adaptive depth. Each stage’s stability is the next one’s precondition.",dict(size=20,color=INK))]])
mcard(25.9,3.9,TEAL,"WHY IT CAN EXIST AT ALL",
     [[("Latent reasoning is trainable ",dict(size=20,color=INK)),("only if",dict(size=20,bold=True,color=TEAL)),
       (" collapse and instability are held off ",dict(size=20,color=INK)),("at once",dict(size=20,bold=True,color=INK)),
       (" — which the anchor + norm-bounded loop + differentiable memory do jointly.",dict(size=20,color=INK))]],tsize=22)

# ---------- THE BET (pale band + rules) ----------
rect(CX,30.55,CW,3.05,fill=PALE[NAVY],lc=SOFT,lw=1.4,radius=0.03)
line(CX+1.4,31.0,CX+CW-1.4,31.0,color=NAVY,w=2.6)
T(CX,31.18,CW,0.7,[[("THE BET",dict(size=24,bold=True,color=TEAL,font=SERIF,spacing=2.6))]],align=PP_ALIGN.CENTER,sa=0)
T(CX+0.6,31.95,CW-1.2,1.4,[[("Reason in a compact latent space over whole thoughts — emit tokens only to ",dict(size=27,color=NAVY,font=SERIF)),
    ("speak",dict(size=27,italic=True,bold=True,color=TEAL,font=SERIF)),(", never to ",dict(size=27,color=NAVY,font=SERIF)),
    ("think",dict(size=27,italic=True,bold=True,color=AMBER,font=SERIF)),(".",dict(size=27,color=NAVY,font=SERIF))]],align=PP_ALIGN.CENTER,line=1.1,sa=0)
line(CX+1.4,33.2,CX+CW-1.4,33.2,color=NAVY,w=2.6)
T(CX,33.95,CW,0.7,[[("Built from four ideas — JEPA-Reasoner · HRM-Text · Thought Gestalt · Parcae      ·      Code + full spec:  ",dict(size=16,color=MUTE)),
    ("github.com/danerjin/hlra",dict(size=16,bold=True,color=TEAL))]],align=PP_ALIGN.CENTER,sa=0)

import sys
OUT=sys.argv[1] if len(sys.argv)>1 else os.path.join(os.path.dirname(os.path.abspath(__file__)),"HLRA_Poster.pptx")
prs.save(OUT)
print("saved %s  (48 x 36 in, content off the hinges at x=12/36)"%OUT)
for png,slot in (("fig_loss.png","Loss vs. steps"),("fig_arc.png","ARC-C vs. FLOPs")):
    print("  [%s] %s"%("figure" if os.path.exists(os.path.join(FIGS,png)) else "placeholder",slot))
