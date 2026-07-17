#!/usr/bin/env python
"""Tile the 48x36 HLRA poster into 16 letter pages (4x4), keeping content off
the tri-fold hinges (x=12, x=36). Tile columns map to the panels:
  col1 = left panel, col2+col3 = center panel, col4 = right panel;
the hinge gaps (11.5–13, 35–36.5) and edge/top/bottom margins print blank.

USAGE (re-tile after editing the master):
  1. In PowerPoint: File > Export (or Save As) > PDF  ->  HLRA_Poster.pdf
  2. .venv/bin/python tile.py HLRA_Poster.pdf
     (or pass a PNG you exported:  .venv/bin/python tile.py poster.png)
Needs: pdftoppm (brew install poppler) only when the input is a PDF; else just PIL.
"""
import sys, subprocess, os, io
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DPI=200
MARGIN_MM=4.0                 # printer's unprintable border on each edge
# tile bands in inches on the 48x36 canvas (each 11 x 8.5); off the hinges
XB=[(0.5,11.5),(13.0,24.0),(24.0,35.0),(36.5,47.5)]
YB=[(1.0,9.5),(9.5,18.0),(18.0,26.5),(26.5,35.0)]
PW,PH=11.0,8.5                # letter page (landscape)

def load(src):
    if src.lower().endswith(".pdf"):
        subprocess.run(["pdftoppm","-png","-r",str(DPI),src,"_hlra_page"],check=True)
        img="_hlra_page-1.png"; im=Image.open(img).convert("RGB"); os.remove(img); return im
    return Image.open(src).convert("RGB")

src=sys.argv[1] if len(sys.argv)>1 else "HLRA_Poster.pdf"
im=load(src); Wp,Hp=im.size
sx,sy=Wp/48.0,Hp/36.0        # px per inch (from whatever the source resolution is)

mm=MARGIN_MM/25.4
availW,availH=PW-2*mm,PH-2*mm          # printable area
prs=Presentation(); prs.slide_width=Inches(PW); prs.slide_height=Inches(PH)
blank=prs.slide_layouts[6]
for r,(y0,y1) in enumerate(YB):
    for c,(x0,x1) in enumerate(XB):
        tile=im.crop((round(x0*sx),round(y0*sy),round(x1*sx),round(y1*sy)))
        buf=io.BytesIO(); tile.save(buf,format="PNG"); buf.seek(0)
        sl=prs.slides.add_slide(blank)
        # scale the whole tile to fit inside the printable area (uniform, aspect kept)
        tw,th=(x1-x0),(y1-y0); sc=min(availW/tw,availH/th)
        dw,dh=tw*sc,th*sc; mx,my=(PW-dw)/2,(PH-dh)/2
        sl.shapes.add_picture(buf,Inches(mx),Inches(my),Inches(dw),Inches(dh))
        # (no corner markers) — slides are ordered row-major: R1C1..R1C4, R2C1..
prs.save("HLRA_Poster_PrintTiles.pptx")
print("wrote HLRA_Poster_PrintTiles.pptx  (16 tiles; %.0fmm safe border, content off the hinges)"%MARGIN_MM)

ref=im.copy(); ref.thumbnail((3400,2600)); ref.save("HLRA_Poster_reference.png")
print("wrote HLRA_Poster_reference.png", ref.size)
